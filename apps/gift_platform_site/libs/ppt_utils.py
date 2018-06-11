from pptx import Presentation
from pptx.util import Cm,Pt
from django.conf import settings
from pptx.enum.text import MSO_AUTO_SIZE
import os


def generate_ppt(product_list, path):
    """
    生成PPT
    :param product_list: 商品列表
    :return:
    """
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for product in product_list:
        slide = prs.slides.add_slide(blank_slide_layout)

        # logo
        # slide.shapes.add_picture(os.path.join(settings.BASE_DIR, product.brand.logo.path), Cm(0), Cm(0), Cm(4.02), Cm(2.61))
        #

        # 主图
        left = Cm(1.59)
        top = Cm(2.98)
        width = Cm(10.92)
        height = Cm(10.92)

        pic = slide.shapes.add_picture(os.path.join(settings.BASE_DIR, product.images.first().productimage.path) , left, top, width, height)
        #

        # 规格图片
        print(len(product_list))
        spec_pics = product.images.all()[1:4]
        spec_img_width = Cm(2.5)
        spec_img_space = Cm(0.31)
        spec_img_top = Cm(14.6)
        spec_img_start_width = Cm(1.59)
        i = 0

        for spec_pic in spec_pics:
            slide.shapes.add_picture(os.path.join(settings.BASE_DIR, spec_pic.productimage.path), Cm(1.59 + i * 2.5 + i * 0.31) , spec_img_top, spec_img_width)
            i = i + 1

        # 描述
        desc_box = slide.shapes.add_textbox(Cm(14.29), Cm(2.78), Cm(10.56), Cm(15.71))
        tf = desc_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "{}".format(product.name)
        p.font.bold = True
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(28)

        if product.brand is not None:
            p = tf.add_paragraph()
            p.font.bold = True
            p.font.size = Pt(15)
            p.text = "品牌：{}".format(product.brand.name)

        if product.model is not None:
            p = tf.add_paragraph()
            p.font.bold = True
            p.font.size = Pt(15)
            p.text = "型号：{}".format(product.model)

        # 价格
        # price_box = slide.shapes.add_textbox(Cm(14.29), Cm(4.11), Cm(6.15), Cm(2.82))
        #tf = price_box.text_frame
        #tf.clear()
        product_items = product.productItems.order_by("price").all()
        start_price = product_items.first().price
        end_price = product_items.last().price

        #p = tf.paragraphs[0]
        p = tf.add_paragraph()
        p.font.name = 'Microsoft YaHei'
        price_content = ""
        if start_price == end_price:
            price_content = "市场价：{} 元".format(start_price)
        else:
            price_content = "市场价：{} - {} 元".format(start_price, end_price)
        p.font.bold = True
        p.font.size = Pt(15)
        price_content = price_content + '\n'

        # 计算供货价
        product_items = product.productItems.order_by("favouredprice").all()
        start_favored_price = product_items.first().favouredprice
        end_favored_price = product_items.last().favouredprice

        if start_favored_price != 0 and end_favored_price != 0:
            if start_price == end_price:
                price_content = price_content + "供货价：{} 元".format(start_favored_price)
            else:
                price_content = price_content + "供货价：{} - {} 元".format(start_favored_price, end_favored_price)
        #
        p.text = price_content
        #

        # 规格添加
        #details_box = slide.shapes.add_textbox(Cm(14.29), Cm(4.70), Cm(6.15), Cm(2.82))
        #tf = details_box.text_frame
        #tf.clear()

        p = tf.add_paragraph()

        # 开始添加规格描述：格式为：颜色：红 黄
        specname_array = product.attributes
        spec_desc_text = ''
        for spec_name in specname_array:
            text_row = '{}:'.format(spec_name)
            for sku in product.productItems.all():
                if sku.attributes.get(spec_name) not in text_row: # 排除掉已经添加的规格
                    text_row = text_row + sku.attributes.get(spec_name) + ' '
            spec_desc_text = spec_desc_text + text_row + '\n'
            #p = tf.add_paragraph()
        p.font.bold = False
        p.font.size = Pt(15)
        p.text = spec_desc_text
        p.line_spacing = 1.5
        #


        #title_box = slide.shapes.add_textbox(Cm(14.29), Cm(8.37), Cm(10.47), Cm(9.9))
        #tf = title_box.text_frame
        #tf.word_wrap = True
        #tf.clear()
        #tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        #p = tf.paragraphs[0]
        p = tf.add_paragraph()
        p.font.name = 'Microsoft YaHei'
        p.text = "产品介绍："
        p.font.bold = True
        p.font.size = Pt(11)
        p = tf.add_paragraph()
        p.font.bold = False
        p.font.size = Pt(11)
        p.font.name = 'Microsoft YaHei'
        if product.simple_description is not None:
            p.text = product.simple_description
        p.line_spacing = 1.5

        """
        desc_box = slide.shapes.add_textbox(Cm(10.58 + 1.32), Cm(1.59 + 1.59 + 7.14 + 0.79), Cm(13.23), Cm(5.29))
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "TODO:待定"
        p.font.bold = False
        p.font.size = Pt(14)
        """

    if not os.path.exists(os.path.dirname(path)):
        os.makedirs(os.path.dirname(path))
    prs.save(path)
